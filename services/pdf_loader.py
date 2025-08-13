# Import necessary libraries
import os
from langchain.schema import Document
from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation

# -------------------------------------------------------
# Function: Split documents into smaller text chunks
# -------------------------------------------------------
def split_text_into_chunks(documents, chunk_size = 1000, chunk_overlap = 200):
    # Create a text splitter that divides text into chunks
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size, 
        chunk_overlap=chunk_overlap
    )
    return text_splitter.split_documents(documents)

# -------------------------------------------------------
# Function: Extract text from different file formats
# -------------------------------------------------------
def extract_text_from_document(file_path):

    # ---------- PDF FILES ----------
    if file_path.lower().endswith('.pdf'):
        # Use PyPDFLoader to read PDF files into Document objects
        loader = PyPDFLoader(file_path)
        documents = loader.load()

    # ---------- POWERPOINT FILES ----------
    elif file_path.lower().endswith(('.pptx', '.pptm')):
        # Load the PowerPoint file
        prs = Presentation(file_path)
        documents = []

        # Loop through each slide
        for i, slide in enumerate(prs.slides, start=1):
            text_runs = []

            # Loop through each shape (text box, title, etc.)
            for shape in slide.shapes:
                # If shape contains text, add it to text_runs
                if hasattr(shape, "text") and shape.text.strip():
                    text_runs.append(shape.text.strip())
                
            # If the slide has text, create a Document object for it
            if text_runs:
                documents.append(Document(
                    page_content="\n".join(text_runs),
                    metadata={"source": str(file_path), "page": i}
                ))

    # ---------- TEXT FILES ----------
    elif file_path.lower().endswith('.txt'):
        # Get the file name only (without full path)
        source_name = os.path.basename(file_path)

        # Read the entire text file
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read().strip()

        # Store entire file as one Document
        documents = [
            Document(page_content=text, metadata={"source": source_name, "page": 1})
        ]

    # ---------- UNSUPPORTED FILE TYPES ----------
    else:
        raise ValueError("Unsupported file type. Please upload a PDF, PPTX, PPTM, or TXT file.")

    return documents
